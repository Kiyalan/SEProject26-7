"""
Patch key claims in 技术原型迭代评审.pptx to match SEProject26-7 reality.
Creates a .bak backup beside the original file.
"""
from __future__ import annotations

import re
import shutil
from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

PPTX = Path(r"C:\Users\onayi\Desktop") / "小学期项目" / "技术原型迭代评审.pptx"
BACKUP = PPTX.with_suffix(".pptx.bak")

# Exact full-string replacements (shape/paragraph text after strip)
EXACT = {
    # Slide 7 highlights
    "🔐 多用户全隔离：完整 JWT 鉴权体系，用户级、仓库级双重数据隔离，支持云部署多租户场景":
        "🔐 多用户隔离：JWT 鉴权 + owner_login/仓库权限校验；适合本地/演示多用户，非完整 SaaS 多租户",
    "🤖 FAQ 自动生成：基于知识库内容自动聚类生成 FAQ，而非纯手动维护，更贴合长期记忆的业务定位":
        "🤖 FAQ 自动生成：基于 GraphRAG 证据按主题聚类生成 FAQ，并可导出；手工增删改仍可继续补齐",
    "✨ 增强型 RAG 检索：在基础向量检索之上，新增 AST 语法树解析与 GraphRAG 图关系检索，代码场景准确率更高":
        "✨ 增强型 RAG：CodeWiki AST 解析 + GraphRAG 图检索；向量 Embedding 可选开启，默认图扩展+全文检索",
    # Slide 14 stack
    "Java 17 + Spring Boot 2.7 + MyBatis-Plus":
        "Java 21 + Spring Boot 3.3 + Spring Data JPA / JdbcTemplate",
    "MySQL 8.0 + Milvus 向量数据库":
        "H2（业务投影）+ PostgreSQL/pgvector（CodeWiki 图与检索）",
    "通义千问 API + Embedding 模型":
        "可配置 LLM API（如 OpenRouter）+ 可选 Embedding",
    "GitHub REST API + Webhook":
        "GitHub REST API（OAuth / Issues / Actions；Webhook 待下迭代）",
    "Docker + Linux 云服务器":
        "Docker Compose（postgres + codewiki）+ 本地/云主机部署",
    "工程化能力：Docker 容器化部署、数据库版本迁移、环境变量自动加载":
        "工程化能力：Docker Compose 知识引擎、schema.sql 初始化、环境变量/.env 加载",
    # Slide 15 algorithm pipeline
    "代码采集 → AST语法解析 → 语义切分 → 向量化 + 图关系构建 → 混合相似度检索 → LLM生成回答":
        "JGit 同步 → CodeWiki AST 分析 → GraphRAG 构建 → 图/全文(/可选向量)检索 → LLM 生成回答",
    "3.混合检索策略：向量语义 + 关键词 + 图关系三重匹配，检索精准度更高":
        "3.混合检索策略：图关系 + 全文检索（+ 可选向量）匹配，适配代码跨文件关联",
    "4.可迭代优化：切分策略、相似度阈值、Prompt 均可迭代调优，适配不同代码仓库":
        "4.可迭代优化：检索意图路由、Prompt、Embedding 开关均可调优，适配不同仓库",
    # Slide 16
    "3.最终采用「向量语义 + 关键词匹配 + 图关系关联」三重混合检索策略":
        "3.当前落地「图关系 + 全文检索（可选向量）」混合策略，由 CodeWiki GraphRAG 承载",
    "2.可迭代优化：切分策略、相似度阈值、权重配比均可调优":
        "2.可迭代优化：意图路由、证据条数、LLM Prompt 与 Embedding 开关可调优",
    # Slide 19/20 metrics - soften
    "算法效果：AST+GraphRAG 增强方案，代码场景检索召回率相比纯文本 RAG 提升约 30%":
        "算法效果：AST+GraphRAG 方案已跑通；量化召回对比待用评测集补齐（勿写未实测百分比）",
    "算法效果：AST+GraphRAG 增强方案，代码场景检索召回率相比纯文本基线提升约 30%":
        "算法效果：AST+GraphRAG 方案已跑通；量化召回对比待用评测集补齐（勿写未实测百分比）",
    # Slide 22 quality table cells
    "100%（异常捕获 + 断点续建 + 并发锁）":
        "显著提升（异步任务 + 异常兜底；大仓分片/断点续建待下迭代）",
    "全链路 JWT + 租户级双维度隔离":
        "JWT + owner_login/仓库授权隔离（演示级多用户）",
    "Docker 容器化 + 多环境配置，运行完全一致":
        "Docker Compose 知识引擎 + 开发配置；云生产压测待验证",
    # Slide 25 risk mitigation
    "本轮已落地缓解：1.引入 AST 抽象语法树按代码结构切分片段，保留代码原生语义； 2.新增 GraphRAG 构建代码调用关系图谱，补充跨文件关联检索； 3.搭建「向量语义 + 关键词 + 图关系」三重混合检索策略，召回率提升超 25%":
        "本轮已落地缓解：1.CodeWiki AST 保留代码结构语义；2.GraphRAG 构建调用/依赖图；3.图扩展+全文（可选向量）检索。量化召回提升待评测集验证",
    "后续迭代优化：引入异步任务队列 + 分片并行构建，前端展示实时构建进度":
        "后续迭代优化：强化分片并行/断点续建；前端已有进度与任务日志，可持续增强",
    # Slide 26 tenant
    "本轮已落地缓解：1.搭建全局 JWT 鉴权体系，所有接口统一拦截校验用户身份； 2.全数据表新增user_id租户字段，查询强制绑定当前用户做数据过滤； 3.区分普通用户 / 管理员的接口权限粒度":
        "本轮已落地缓解：1.JWT 鉴权；2.核心表 owner_login + GitHub 仓库授权校验；3.普通用户/管理员接口分离。完整 SaaS 租户模型待加强",
    # Slide 30 improvements
    "✅ 体系补全：落地完整 JWT 鉴权与双维度数据隔离，支持多用户云部署":
        "✅ 体系补全：落地 JWT 鉴权与 owner_login/仓库隔离，支持本地多用户演示",
    "✅ 工程优化：双数据源支持、数据库迁移脚本、环境变量自动加载、容器化部署方案":
        "✅ 工程优化：H2+CodeWiki/pgvector、schema 初始化、环境变量加载、Docker Compose 方案",
    # Slide 31 next iteration - fix contradictory "接入前端"
    "3.功能补全：完善 UC1 仓库同步、UC3 完整问答、UC4 Issue 分析、UC5 Bug 修复全链路":
        "3.功能补全：FAQ 手工增删改、Webhook 同步、评测集与召回度量、Issue/PR 链路加固",
    "4.体验优化：接入前端页面，完成前后端联调，补充邮件通知等辅助功能":
        "4.体验优化：前端已联调；继续完善邮件 SMTP、加载态与管理端体验",
    "1.性能优化：引入异步任务队列，实现大仓库并行分片构建，解决卡死风险":
        "1.性能优化：在现有异步构建之上，推进大仓库分片并行与断点续建",
    # Slide 6 placeholder hint
    "此处需要 4 张真实截图，每张配一句话说明":
        "请替换为 4 张真实截图：①接口 200 ②数据表有数据 ③知识库/问答效果 ④IDEA 包结构",
    # Slide 12 knowledge service duty
    "向量索引构建、长期记忆沉淀":
        "CodeWiki/GraphRAG 索引构建、FAQ 聚类与导出",
    # Slide 5 demo UC6
    "3.FAQ 自动聚类生成与手动维护（UC6）":
        "3.FAQ 自动聚类生成与导出（UC6；手工维护接口可继续补）",
}

# Substring replacements applied when exact match fails (order matters: longer first)
PARTIAL = [
    ("召回率相比纯文本 RAG 提升约 30%", "方案已跑通；量化召回对比待评测集补齐"),
    ("召回率相比纯文本基线提升约 30%", "方案已跑通；量化召回对比待评测集补齐"),
    ("召回率提升超 25%", "量化召回提升待评测集验证"),
    ("提升约 30%", "（待评测）"),
    ("提升约 18%", "（待评测）"),
    ("Java 17 + Spring Boot 2.7 + MyBatis-Plus",
     "Java 21 + Spring Boot 3.3 + JPA/JdbcTemplate"),
    ("MySQL 8.0 + Milvus 向量数据库",
     "H2 + PostgreSQL/pgvector（CodeWiki）"),
    ("通义千问 API + Embedding 模型",
     "可配置 LLM API + 可选 Embedding"),
    ("GitHub REST API + Webhook",
     "GitHub REST API（Webhook 待下迭代）"),
    ("支持云部署多租户场景",
     "支持本地/演示多用户隔离"),
    ("支持多用户云部署",
     "支持本地多用户演示"),
    ("断点续建 + 并发锁",
     "异步任务 + 异常兜底"),
    ("user_id租户字段",
     "owner_login 归属字段"),
    ("三重混合检索",
     "图+全文（可选向量）检索"),
    ("向量语义 + 关键词 + 图关系",
     "图关系 + 全文（可选向量）"),
    ("向量语义 + 关键词匹配 + 图关系关联",
     "图关系 + 全文检索（可选向量）"),
    ("接入前端页面，完成前后端联调，补充邮件通知等辅助功能",
     "前端已联调；继续完善邮件 SMTP 与体验细节"),
    ("长期记忆沉淀",
     "FAQ/知识库持久化"),
    ("Milvus", "pgvector/CodeWiki"),
]


def iter_shapes(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)


def replace_in_textframe(tf) -> int:
    changed = 0
    # Prefer whole-frame text replace when the entire frame matches
    full = (tf.text or "").strip()
    if full in EXACT:
        new = EXACT[full]
        if tf.paragraphs:
            # Keep first run formatting; clear other paragraphs' text lightly
            p0 = tf.paragraphs[0]
            if p0.runs:
                p0.runs[0].text = new
                for r in p0.runs[1:]:
                    r.text = ""
            else:
                p0.text = new
            for p in tf.paragraphs[1:]:
                for r in p.runs:
                    r.text = ""
                if not p.runs:
                    p.text = ""
        changed += 1
        return changed

    for p in tf.paragraphs:
        raw = p.text or ""
        if not raw.strip():
            continue
        new = raw
        stripped = raw.strip()
        if stripped in EXACT:
            # preserve leading/trailing whitespace of paragraph
            lead = raw[: len(raw) - len(raw.lstrip())]
            trail = raw[len(raw.rstrip()) :]
            new = lead + EXACT[stripped] + trail
        else:
            for old, repl in PARTIAL:
                if old in new:
                    new = new.replace(old, repl)
        if new != raw:
            if p.runs:
                p.runs[0].text = new
                for r in p.runs[1:]:
                    r.text = ""
            else:
                p.text = new
            changed += 1
    return changed


def replace_in_table(table) -> int:
    changed = 0
    for row in table.rows:
        for cell in row.cells:
            if cell.text_frame is not None:
                changed += replace_in_textframe(cell.text_frame)
    return changed


def main() -> None:
    if not PPTX.exists():
        raise SystemExit(f"missing: {PPTX}")
    if not BACKUP.exists():
        shutil.copy2(PPTX, BACKUP)
        print(f"backup -> {BACKUP}")
    else:
        print(f"backup exists: {BACKUP}")

    prs = Presentation(str(PPTX))
    total = 0
    for i, slide in enumerate(prs.slides, 1):
        slide_changed = 0
        for shape in iter_shapes(slide.shapes):
            if shape.has_text_frame:
                c = replace_in_textframe(shape.text_frame)
                slide_changed += c
            if shape.has_table:
                c = replace_in_table(shape.table)
                slide_changed += c
        if slide_changed:
            print(f"slide {i}: {slide_changed} text unit(s) updated")
            total += slide_changed

    prs.save(str(PPTX))
    print(f"saved {PPTX} total_units={total}")


if __name__ == "__main__":
    main()
